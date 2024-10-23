import os
import io
import csv
import docx
import pdfplumber
from pptx import Presentation
from PIL import Image

class UniversalDataExtractor:
    """
    UniversalDataExtractor class handles the extraction of text, tables, images, metadata, and links
    from various file formats: PDF, DOCX, and PPTX.
    """

    def __init__(self, loader):
        """
        Initializes the UniversalDataExtractor with a file loader.
        
        :param loader: An instance of a file loader that provides the file to be processed.
        """
        self.file_loader = loader
        self.content = self.file_loader.load_file()  # Load the file's content using the loader
        self.file_type = os.path.splitext(loader.file_path)[1].lower()  # Get the file extension
        
        # Open the appropriate file type
        if self.file_type == '.pdf':
            self.pdf = pdfplumber.open(self.file_loader.file_path)
        elif self.file_type == '.docx':
            self.doc = docx.Document(self.file_loader.file_path)
        elif self.file_type == '.pptx':
            self.prs = Presentation(self.file_loader.file_path)

    def extract_text(self):
        """
        Extracts the text from the file based on its type (PDF, DOCX, PPTX).
        
        :return: Extracted text as a string.
        """
        extracted_data = {'text': '', 'metadata': {}}

        # Extract text based on the file type
        if self.file_type == '.pdf':
            return "\n".join(page.extract_text() or "" for page in self.pdf.pages).strip()
        elif self.file_type == '.docx':
            return "\n".join(para.text for para in self.doc.paragraphs).strip()
        elif self.file_type == '.pptx':
            text = ''
            for slide in self.prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()

        # If no text is found, return an empty string
        return ""

    def extract_tables(self):
        """
        Extracts tables from the file based on its type (PDF, DOCX, PPTX).
        
        :return: A list of extracted tables.
        """
        tables = []
        
        # Extract tables based on the file type
        if self.file_type == '.pdf':
            for page in self.pdf.pages:
                tables.extend(page.extract_tables() or [])  # Extract tables from each PDF page
        elif self.file_type == '.docx':
            tables = [[self._extract_table_row(row) for row in table.rows] for table in self.doc.tables]
        elif self.file_type == '.pptx':
            for slide in self.prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        tables.append([self._extract_table_row(row) for row in shape.table.rows])

        return tables
    
    def _extract_table_row(self, row):
        """
        Helper method to extract text from a table row.
        
        :param row: A table row object.
        :return: A list of cell text in the row.
        """
        return [cell.text.strip() for cell in row.cells]
    
    def extract_images(self):
        """
        Extracts images from the file based on its type (PDF, DOCX, PPTX).
        
        :return: A list of paths to the saved images.
        """
        images = []
        
        # Extract images based on the file type
        if self.file_type == ".pdf":
            for page_number, page in enumerate(self.pdf.pages):
                for img_index, img in enumerate(page.images):
                    if 'stream' in img:
                        img_data = img['stream'].get_rawdata()  # Get image data
                        images.append(self.save_image(img_data, img_index + 1, ".pdf", page_number))  # Save image
        elif self.file_type == ".docx":
            for rel in self.doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob  # Get image data
                    images.append(self.save_image(img_data, len(images) + 1, ".docx"))  # Save image
        elif self.file_type == ".pptx":
            for slide_number, slide in enumerate(self.prs.slides):
                for shape_index, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # Check if the shape contains an image
                        img_data = shape.image.blob  # Get image data
                        images.append(self.save_image(img_data, shape_index + 1, ".pptx", slide_number))  # Save image
        return images
    
    def extract_metadata(self):
        """
        Extracts metadata from the file based on its type (PDF, DOCX, PPTX).
        
        :return: Metadata as a dictionary.
        """
        if self.file_type == ".pdf":
            return self.pdf.metadata
        elif self.file_type == ".docx":
            return self.doc.core_properties
        elif self.file_type == ".pptx":
            return self.prs.core_properties
        
        return {}
    
    def extract_links(self):
        """
        Extracts hyperlinks from the file based on its type (PDF, DOCX, PPTX).
        
        :return: A list of extracted links.
        """
        links = []
        
        # Extract links based on the file type
        if self.file_type == ".pdf":
            for page in self.pdf.pages:
                links.extend(annot.get("uri") for annot in getattr(page, 'annots', []) if annot.get("uri"))
        elif self.file_type == ".docx":
            for rel in self.doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    links.append(rel.target_ref)
        elif self.file_type == ".pptx":
            for slide in self.prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.hyperlink:
                                    links.append(run.hyperlink.address)
        
        return links
    
    def get_file_name(self):
        """
        Returns the base name of the file being processed.
        
        :return: The name of the file without the directory path.
        """
        return os.path.basename(self.file_loader.file_path)
    
    def save_image(self, img_data, index, file_ext, page_number=None):
        """
        Saves an image to the output directory and returns the path to the saved image.
        
        :param img_data: Raw image data.
        :param index: Index to append to the image filename.
        :param file_ext: File extension of the original file.
        :param page_number: Page number if the image is from a PDF.
        :return: The path to the saved image file.
        """
        img_filename = f"{self.get_file_name().replace(file_ext, '')}_img_{index}.png"
        if page_number is not None:
            img_filename = f"{self.get_file_name().replace(file_ext, '')}_page_{page_number + 1}_img_{index}.png"
        
        img_path = os.path.join('output', img_filename)  # Set the output path for the image
        image = Image.open(io.BytesIO(img_data))  # Open image from raw data
        image.save(img_path)  # Save the image
        return img_path
    
    def close(self):
        """
        Closes the open file resources to free up memory.
        """
        if self.file_type == '.pdf':
            self.pdf.close()  # Close the PDF file when done
